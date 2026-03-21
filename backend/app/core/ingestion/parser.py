"""
Hybrid document parser — upgraded from previous Flask project.

Strategy (same resilient multi-lib approach, now async-ready):
  PDF  : PyMuPDF primary → pdfplumber fallback → PyPDF2 last resort
  DOCX : python-docx
  PPTX : python-pptx
  TXT/MD: plain read

Returns a ParsedDocument with raw text, heading list,
page texts, and raw image bytes for further processing.
"""
from __future__ import annotations
import io
import os
import re
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
from loguru import logger


@dataclass
class RawImage:
    """Image extracted from a document before embedding."""
    id: str = field(default_factory=lambda: str(uuid.uuid4()))
    page_number: int = 0
    image_bytes: bytes = b""
    image_ext: str = "png"
    width: int = 0
    height: int = 0
    # Spatial context (from prev project — Y-coord association)
    bbox_y0: float = 0.0
    bbox_y1: float = 0.0
    spatial_nearest_heading: Optional[str] = None
    spatial_distance: float = float("inf")
    # Caption / fig label (detected during parsing)
    caption: Optional[str] = None
    fig_label: Optional[str] = None


@dataclass
class ParsedDocument:
    """Output of the parser — input to tree_builder."""
    doc_id: str
    filename: str
    mime_type: str
    full_text: str = ""
    # Per-page text (used for spatial image association)
    pages: list[str] = field(default_factory=list)
    # Extracted headings in order: (heading_text, level, char_offset)
    headings: list[tuple[str, int, int]] = field(default_factory=list)
    # Raw images before CLIP embedding
    images: list[RawImage] = field(default_factory=list)
    total_pages: int = 0
    title: str = ""


# ── Caption / fig-label patterns ────────────────────────────────────────
FIG_LABEL_RE = re.compile(
    r"\b(Fig(?:ure)?[\s.]*(\d+[a-zA-Z]?|[A-Z]))\b",
    re.IGNORECASE,
)
CAPTION_RE = re.compile(
    r"^(Fig(?:ure)?[\s.]*[\d\w]+[:\s]+.{5,120})$",
    re.IGNORECASE | re.MULTILINE,
)


def _extract_fig_label(text: str) -> Optional[str]:
    m = FIG_LABEL_RE.match(text.strip())
    return m.group(1) if m else None


def _find_caption_near(text: str, approx_pos: int, window: int = 400) -> Optional[str]:
    """Search for a caption string near a character position."""
    start = max(0, approx_pos - window)
    end = min(len(text), approx_pos + window)
    snippet = text[start:end]
    m = CAPTION_RE.search(snippet)
    return m.group(1).strip() if m else None


# ══════════════════════════════════════════════════════════
# PDF parser
# ══════════════════════════════════════════════════════════
def _parse_pdf(filepath: str, doc_id: str) -> ParsedDocument:
    filename = Path(filepath).name
    doc = ParsedDocument(doc_id=doc_id, filename=filename, mime_type="application/pdf")

    # ── Primary: PyMuPDF ─────────────────────────────────────────────────
    try:
        import fitz  # PyMuPDF
        pdf = fitz.open(filepath)
        doc.total_pages = len(pdf)
        full_parts: list[str] = []

        for page_num, page in enumerate(pdf):
            # Text
            page_text = page.get_text("text") or ""
            doc.pages.append(page_text)
            full_parts.append(page_text)

            # Headings from text blocks (from prev project pattern, enhanced)
            blocks = page.get_text("dict").get("blocks", [])
            for block in blocks:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue
                    line_text = "".join(s.get("text", "") for s in spans).strip()
                    if not line_text:
                        continue
                    # Detect heading by font size + boldness
                    avg_size = sum(s.get("size", 0) for s in spans) / len(spans)
                    is_bold = any("bold" in s.get("font", "").lower() for s in spans)
                    char_offset = len("\n".join(full_parts))
                    if avg_size >= 13 or (avg_size >= 11 and is_bold):
                        level = 1 if avg_size >= 16 else (2 if avg_size >= 13 else 3)
                        doc.headings.append((line_text, level, char_offset))

            # Images with spatial Y-coord (prev project innovation, enhanced)
            img_list = page.get_images(full=True)
            page_headings_spatial = []

            for block in blocks:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    line_text = "".join(
                        s.get("text", "") for s in line.get("spans", [])
                    ).strip()
                    if line_text:
                        bbox = line.get("bbox", [0, 0, 0, 0])
                        page_headings_spatial.append({
                            "text": line_text,
                            "y0": bbox[1],
                        })

            for img_info in img_list:
                xref = img_info[0]
                try:
                    base_img = pdf.extract_image(xref)
                    img_bytes = base_img.get("image", b"")
                    if not img_bytes or len(img_bytes) < 500:  # skip tiny images
                        continue
                    img_rect = page.get_image_bbox(img_info)
                    img_y = img_rect.y0 if img_rect else 0.0

                    # Find nearest heading by Y distance (prev project core idea)
                    nearest_heading = None
                    min_dist = float("inf")
                    for h in page_headings_spatial:
                        dist = abs(h["y0"] - img_y)
                        if dist < min_dist:
                            min_dist = dist
                            nearest_heading = h["text"]

                    # Find caption near image position
                    approx_char = len("\n".join(full_parts))
                    caption = _find_caption_near(page_text, 0)
                    fig_label = _extract_fig_label(caption) if caption else None

                    raw_img = RawImage(
                        page_number=page_num + 1,
                        image_bytes=img_bytes,
                        image_ext=base_img.get("ext", "png"),
                        width=base_img.get("width", 0),
                        height=base_img.get("height", 0),
                        bbox_y0=img_y,
                        spatial_nearest_heading=nearest_heading,
                        spatial_distance=min_dist,
                        caption=caption,
                        fig_label=fig_label,
                    )
                    doc.images.append(raw_img)
                except Exception as e:
                    logger.debug(f"Skipping image xref={xref}: {e}")

        doc.full_text = "\n".join(full_parts)
        doc.title = _extract_title(doc.full_text, filename)
        pdf.close()
        logger.info(f"PDF parsed via PyMuPDF: {len(doc.pages)} pages, "
                    f"{len(doc.headings)} headings, {len(doc.images)} images")
        return doc

    except Exception as e:
        logger.warning(f"PyMuPDF failed: {e} — trying pdfplumber")

    # ── Fallback 1: pdfplumber ────────────────────────────────────────────
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            doc.total_pages = len(pdf.pages)
            parts = []
            for page in pdf.pages:
                t = page.extract_text() or ""
                doc.pages.append(t)
                parts.append(t)
            doc.full_text = "\n".join(parts)
            doc.title = _extract_title(doc.full_text, filename)
        logger.info(f"PDF parsed via pdfplumber: {doc.total_pages} pages")
        return doc
    except Exception as e:
        logger.warning(f"pdfplumber failed: {e} — trying PyPDF2")

    # ── Fallback 2: PyPDF2 ───────────────────────────────────────────────
    try:
        import PyPDF2
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            doc.total_pages = len(reader.pages)
            parts = []
            for page in reader.pages:
                t = page.extract_text() or ""
                doc.pages.append(t)
                parts.append(t)
            doc.full_text = "\n".join(parts)
            doc.title = _extract_title(doc.full_text, filename)
        logger.info(f"PDF parsed via PyPDF2: {doc.total_pages} pages")
        return doc
    except Exception as e:
        logger.error(f"All PDF parsers failed: {e}")
        doc.full_text = ""
        return doc


# ══════════════════════════════════════════════════════════
# DOCX parser
# ══════════════════════════════════════════════════════════
def _parse_docx(filepath: str, doc_id: str) -> ParsedDocument:
    from docx import Document as DocxDocument
    filename = Path(filepath).name
    doc = ParsedDocument(
        doc_id=doc_id, filename=filename,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    try:
        docx = DocxDocument(filepath)
        parts = []
        char_offset = 0
        for para in docx.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if style.startswith("Heading"):
                try:
                    level = int(style.split()[-1])
                except ValueError:
                    level = 1
                doc.headings.append((text, level, char_offset))
            parts.append(text)
            char_offset += len(text) + 1

        doc.full_text = "\n".join(parts)
        doc.total_pages = 1
        doc.title = _extract_title(doc.full_text, filename)
        logger.info(f"DOCX parsed: {len(doc.headings)} headings")
    except Exception as e:
        logger.error(f"DOCX parse error: {e}")
    return doc


# ══════════════════════════════════════════════════════════
# PPTX parser
# ══════════════════════════════════════════════════════════
def _parse_pptx(filepath: str, doc_id: str) -> ParsedDocument:
    from pptx import Presentation
    filename = Path(filepath).name
    doc = ParsedDocument(
        doc_id=doc_id, filename=filename,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    try:
        prs = Presentation(filepath)
        doc.total_pages = len(prs.slides)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            for shape in slide.shapes:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue
                text = shape.text.strip()
                slide_text_parts.append(text)
                # Title placeholders → H1 headings
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if shape.placeholder_format and shape.placeholder_format.idx == 0:
                        char_offset = len("\n".join(parts))
                        doc.headings.append((text, 1, char_offset))

                # Extract images from shapes
                if hasattr(shape, "image") and shape.image:
                    try:
                        img_bytes = shape.image.blob
                        if img_bytes and len(img_bytes) > 500:
                            raw_img = RawImage(
                                page_number=slide_num,
                                image_bytes=img_bytes,
                                image_ext=shape.image.ext or "png",
                                spatial_nearest_heading=doc.headings[-1][0] if doc.headings else None,
                            )
                            doc.images.append(raw_img)
                    except Exception:
                        pass

            slide_text = "\n".join(slide_text_parts)
            doc.pages.append(slide_text)
            parts.append(slide_text)

        doc.full_text = "\n".join(parts)
        doc.title = _extract_title(doc.full_text, filename)
        logger.info(f"PPTX parsed: {doc.total_pages} slides, {len(doc.headings)} headings, "
                    f"{len(doc.images)} images")
    except Exception as e:
        logger.error(f"PPTX parse error: {e}")
    return doc


# ══════════════════════════════════════════════════════════
# TXT / MD parser
# ══════════════════════════════════════════════════════════
def _parse_text(filepath: str, doc_id: str, mime_type: str) -> ParsedDocument:
    filename = Path(filepath).name
    doc = ParsedDocument(doc_id=doc_id, filename=filename, mime_type=mime_type)
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        doc.full_text = text
        doc.pages = [text]
        doc.total_pages = 1
        doc.title = _extract_title(text, filename)
        # Detect markdown headings
        char_offset = 0
        for line in text.split("\n"):
            m = re.match(r"^(#{1,6})\s+(.+)$", line)
            if m:
                level = len(m.group(1))
                doc.headings.append((m.group(2).strip(), level, char_offset))
            char_offset += len(line) + 1
        logger.info(f"Text parsed: {len(doc.headings)} headings")
    except Exception as e:
        logger.error(f"Text parse error: {e}")
    return doc


# ══════════════════════════════════════════════════════════
# Public entry point
# ══════════════════════════════════════════════════════════
def parse_document(filepath: str, doc_id: str) -> ParsedDocument:
    """
    Route to the correct parser based on file extension.
    Falls back gracefully at every level.
    """
    ext = Path(filepath).suffix.lower().lstrip(".")
    logger.info(f"Parsing document: {filepath} (ext={ext})")

    if ext == "pdf":
        return _parse_pdf(filepath, doc_id)
    elif ext == "docx":
        return _parse_docx(filepath, doc_id)
    elif ext == "pptx":
        return _parse_pptx(filepath, doc_id)
    elif ext in ("txt", "md"):
        mime = "text/markdown" if ext == "md" else "text/plain"
        return _parse_text(filepath, doc_id, mime)
    else:
        logger.error(f"Unsupported extension: {ext}")
        return ParsedDocument(doc_id=doc_id, filename=Path(filepath).name, mime_type="unknown")


def _extract_title(text: str, fallback: str) -> str:
    """Best-effort title extraction from first non-empty line."""
    for line in text.split("\n")[:10]:
        line = line.strip()
        if line and len(line) > 3:
            return line[:120]
    return Path(fallback).stem
